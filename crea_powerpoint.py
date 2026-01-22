"""
Script per creare un PowerPoint con i grafici delle prenotazioni.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Directory del progetto
BASE_DIR = os.path.dirname(__file__)

# Colori tema
COLORE_PRIMARIO = RGBColor(46, 134, 171)   # #2E86AB
COLORE_SCURO = RGBColor(33, 37, 41)        # #212529
COLORE_BIANCO = RGBColor(255, 255, 255)


def aggiungi_slide_titolo(prs):
    """Crea la slide di titolo."""
    slide_layout = prs.slide_layouts[6]  # Layout vuoto
    slide = prs.slides.add_slide(slide_layout)

    # Sfondo colorato
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORE_PRIMARIO
    background.line.fill.background()

    # Titolo principale
    titolo = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = titolo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Analisi Prenotazioni Ristorante"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORE_BIANCO
    p.alignment = PP_ALIGN.CENTER

    # Sottotitolo
    sottotitolo = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sottotitolo.text_frame
    p = tf.paragraphs[0]
    p.text = "Report statistico delle prenotazioni"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORE_BIANCO
    p.alignment = PP_ALIGN.CENTER

    return slide


def aggiungi_slide_grafico(prs, titolo_text, immagine_path, descrizione):
    """Crea una slide con grafico e descrizione."""
    slide_layout = prs.slide_layouts[6]  # Layout vuoto
    slide = prs.slides.add_slide(slide_layout)

    # Titolo slide
    titolo = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = titolo.text_frame
    p = tf.paragraphs[0]
    p.text = titolo_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORE_SCURO
    p.alignment = PP_ALIGN.LEFT

    # Linea decorativa sotto il titolo
    linea = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Inches(0.05)
    )
    linea.fill.solid()
    linea.fill.fore_color.rgb = COLORE_PRIMARIO
    linea.line.fill.background()

    # Immagine del grafico
    if os.path.exists(immagine_path):
        slide.shapes.add_picture(immagine_path, Inches(0.3), Inches(1.3), width=Inches(9.4))

    # Box descrizione in basso
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.6))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = descrizione
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    return slide


def aggiungi_slide_conclusioni(prs):
    """Crea la slide delle conclusioni."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Titolo
    titolo = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = titolo.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusioni e Insight"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORE_SCURO

    # Linea decorativa
    linea = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.05), Inches(2), Inches(0.05)
    )
    linea.fill.solid()
    linea.fill.fore_color.rgb = COLORE_PRIMARIO
    linea.line.fill.background()

    # Contenuto conclusioni
    conclusioni = [
        ("Fasce orarie di punta", "Le ore 20:00-22:00 mostrano la massima occupazione. Considerare personale extra in questi orari."),
        ("Giorni più richiesti", "Il weekend (Venerdì-Sabato) registra il picco di prenotazioni. Ottimizzare la gestione tavoli."),
        ("Trend positivo", "L'andamento degli ultimi 3 mesi mostra una crescita costante delle prenotazioni.")
    ]

    y_pos = 1.5
    for titolo_punto, testo in conclusioni:
        # Box colorato per ogni punto
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(9), Inches(1.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(240, 248, 255)
        box.line.color.rgb = COLORE_PRIMARIO

        # Titolo punto
        txt = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.15), Inches(8.4), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = titolo_punto
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORE_PRIMARIO

        # Testo punto
        txt2 = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.6), Inches(8.4), Inches(0.7))
        tf2 = txt2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = testo
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORE_SCURO

        y_pos += 1.6

    return slide


def main():
    """Crea il PowerPoint completo."""
    print("Creazione PowerPoint in corso...")

    # Crea presentazione
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Titolo
    aggiungi_slide_titolo(prs)
    print("  - Slide titolo creata")

    # Slide 2: Occupazione per fascia oraria
    aggiungi_slide_grafico(
        prs,
        "Occupazione Media per Fascia Oraria",
        os.path.join(BASE_DIR, "grafico_occupazione_oraria.png"),
        "Percentuale media di tavoli occupati durante le diverse fasce orarie di servizio"
    )
    print("  - Slide occupazione oraria creata")

    # Slide 3: Giorni della settimana
    aggiungi_slide_grafico(
        prs,
        "Prenotazioni per Giorno della Settimana",
        os.path.join(BASE_DIR, "grafico_giorni_settimana.png"),
        "Distribuzione delle prenotazioni nei diversi giorni della settimana"
    )
    print("  - Slide giorni settimana creata")

    # Slide 4: Tendenze 3 mesi
    aggiungi_slide_grafico(
        prs,
        "Tendenze Prenotazioni - Ultimi 3 Mesi",
        os.path.join(BASE_DIR, "grafico_tendenze_3_mesi.png"),
        "Andamento settimanale delle prenotazioni con linea di tendenza"
    )
    print("  - Slide tendenze creata")

    # Slide 5: Conclusioni
    aggiungi_slide_conclusioni(prs)
    print("  - Slide conclusioni creata")

    # Salva il file
    output_path = os.path.join(BASE_DIR, "Report_Prenotazioni.pptx")
    prs.save(output_path)

    print(f"\nPowerPoint creato con successo!")
    print(f"File: {output_path}")


if __name__ == "__main__":
    main()